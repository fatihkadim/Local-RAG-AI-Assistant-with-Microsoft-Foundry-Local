"""
Belge yukleme, akilli hash onbellegi ve parcalama (ingestion) modulu.
documents/ klasorundeki belgeleri okur, parcalara boler,
embedding'lerini hesaplar ve Qdrant vektor veritabanina kaydeder.

Performans Ozellikleri:
    - Rust Parser: 14 format icin yuksek performansli yerel C-ABI ayrıştırma
    - Akilli Onbellek (Incremental Hashing): Degismeyen dosyalari atlar, sadece yeni/guncel dosyalari embed eder
    - Toplu Kayit (Bulk Upsert): Qdrant'a tek tek degil, toplu paketler halinde yukler
    - Optimize Batching: GPU/NPU paralelizasyonu ile hizli embedding uretimi
"""

import os
import glob
import json
import hashlib
import time

from src import config
from src import database
from src import embeddings
from src.telemetry import trace_span, record_ingestion_metrics

# ── Rust Parser Import (opsiyonel) ────────────────────────────

try:
    import rust_parser
    RUST_PARSER_AVAILABLE = True
    print("[PARSER] Rust parser yuklendi (yuksek performans modu)")
except ImportError:
    RUST_PARSER_AVAILABLE = False
    print("[PARSER] Rust parser bulunamadi, Python fallback kullanilacak")

# Onbellek dosya yolu
CACHE_FILE = ".ingestion_cache.json"


# ── Hash ve Onbellek Yonetimi ─────────────────────────────────

def _compute_file_hash(filepath):
    """Bir dosyanin SHA-256 hash ozetini hesaplar."""
    hasher = hashlib.sha256()
    try:
        with open(filepath, "rb") as f:
            for chunk in iter(lambda: f.read(65536), b""):
                hasher.update(chunk)
        return hasher.hexdigest()
    except Exception:
        return None


def _load_cache():
    """Onbellek dosyasini yukler."""
    if os.path.exists(CACHE_FILE):
        try:
            with open(CACHE_FILE, "r", encoding="utf-8") as f:
                return json.load(f)
        except Exception:
            return {}
    return {}


def _save_cache(cache):
    """Onbellek dosyasini kaydeder."""
    try:
        with open(CACHE_FILE, "w", encoding="utf-8") as f:
            json.dump(cache, f, indent=2, ensure_ascii=False)
    except Exception as e:
        print(f"  [UYARI] Onbellek kaydedilemedi: {e}")


# ── Python Fallback Parser'lar ────────────────────────────────

def _parse_with_python(filepath, ext):
    """Python kutuphaneleri ile dosya parse eder (Rust yoksa fallback)."""
    if ext in (".txt", ".md", ".rst"):
        return _parse_plaintext(filepath)
    elif ext == ".pdf":
        return _parse_pdf(filepath)
    elif ext == ".docx":
        return _parse_docx_python(filepath)
    elif ext == ".pptx":
        return _parse_pptx_python(filepath)
    elif ext in (".html", ".htm"):
        return _parse_html_python(filepath)
    elif ext in (".csv", ".tsv"):
        return _parse_csv_python(filepath, ext)
    elif ext in (".xlsx", ".xls"):
        return _parse_xlsx_python(filepath)
    elif ext == ".epub":
        return _parse_epub_python(filepath)
    elif ext in (".json", ".jsonl"):
        return _parse_json_python(filepath, ext)
    else:
        raise ValueError(f"Desteklenmeyen format: {ext}")


def _parse_plaintext(filepath):
    """Duz metin dosyasi okur (UTF-8, fallback latin-1)."""
    try:
        with open(filepath, "r", encoding="utf-8") as f:
            return f.read().strip()
    except UnicodeDecodeError:
        with open(filepath, "r", encoding="latin-1") as f:
            return f.read().strip()


def _parse_pdf(filepath):
    """PDF dosyasi okur (pymupdf)."""
    try:
        import pymupdf
    except ImportError:
        raise ImportError("PDF destegi icin 'pymupdf' gerekli: uv pip install pymupdf")

    doc = pymupdf.open(filepath)
    page_texts = [page.get_text() for page in doc]
    content = "\n\n".join(page_texts).strip()
    if not content:
        raise ValueError(f"PDF metin icermiyor (taratilmis PDF olabilir)")
    return content


def _parse_docx_python(filepath):
    """DOCX dosyasi okur (python-docx fallback)."""
    try:
        from docx import Document
    except ImportError:
        raise ImportError("DOCX destegi icin 'python-docx' gerekli: uv pip install python-docx")

    doc = Document(filepath)
    paragraphs = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
    if not paragraphs:
        raise ValueError("DOCX dosyasindan metin cikarilmadi")
    return "\n\n".join(paragraphs)


def _parse_pptx_python(filepath):
    """PPTX dosyasi okur (python-pptx fallback)."""
    try:
        from pptx import Presentation
    except ImportError:
        raise ImportError("PPTX destegi icin 'python-pptx' gerekli: uv pip install python-pptx")

    prs = Presentation(filepath)
    slides_text = []
    for idx, slide in enumerate(prs.slides, 1):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text.strip())
        if texts:
            slides_text.append(f"[Slide {idx}]\n" + "\n".join(texts))

    if not slides_text:
        raise ValueError("PPTX dosyasindan metin cikarilmadi")
    return "\n\n".join(slides_text)


def _parse_html_python(filepath):
    """HTML dosyasi okur (BeautifulSoup fallback)."""
    try:
        from bs4 import BeautifulSoup
    except ImportError:
        raise ImportError("HTML destegi icin 'beautifulsoup4' gerekli: uv pip install beautifulsoup4")

    with open(filepath, "r", encoding="utf-8", errors="replace") as f:
        html_content = f.read()

    soup = BeautifulSoup(html_content, "html.parser")
    for tag in soup(["script", "style", "nav", "footer", "noscript", "svg", "head"]):
        tag.decompose()

    text = soup.get_text(separator="\n", strip=True)
    if not text:
        raise ValueError("HTML dosyasindan metin cikarilmadi")
    return text


def _parse_csv_python(filepath, ext):
    """CSV/TSV dosyasi okur."""
    import csv as csv_module

    delimiter = "\t" if ext == ".tsv" else ","

    with open(filepath, "r", encoding="utf-8", errors="replace") as f:
        reader = csv_module.DictReader(f, delimiter=delimiter)
        rows = []
        for idx, row in enumerate(reader, 1):
            parts = [f"Satir {idx}:"]
            for key, value in row.items():
                if value and value.strip():
                    parts.append(f"  {key}: {value.strip()}")
            if len(parts) > 1:
                rows.append("\n".join(parts))

    if not rows:
        raise ValueError("CSV/TSV dosyasindan veri cikarilmadi")

    return "\n\n".join(rows)


def _parse_xlsx_python(filepath):
    """Excel dosyasi okur (openpyxl fallback)."""
    try:
        from openpyxl import load_workbook
    except ImportError:
        raise ImportError("Excel destegi icin 'openpyxl' gerekli: uv pip install openpyxl")

    wb = load_workbook(filepath, read_only=True, data_only=True)
    all_text = []

    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        rows = list(ws.iter_rows(values_only=True))
        if not rows:
            continue

        headers = [str(h or "").strip() for h in rows[0]]
        sheet_parts = []

        if len(wb.sheetnames) > 1:
            sheet_parts.append(f"[Sheet: {sheet_name}]")

        for idx, row in enumerate(rows[1:], 1):
            row_parts = [f"Satir {idx}:"]
            for col_idx, value in enumerate(row):
                if value is not None and str(value).strip():
                    header = headers[col_idx] if col_idx < len(headers) else "?"
                    row_parts.append(f"  {header}: {str(value).strip()}")
            if len(row_parts) > 1:
                sheet_parts.append("\n".join(row_parts))

        if sheet_parts:
            all_text.append("\n\n".join(sheet_parts))

    wb.close()
    if not all_text:
        raise ValueError("Excel dosyasindan veri cikarilmadi")
    return "\n\n---\n\n".join(all_text)


def _parse_epub_python(filepath):
    """EPUB dosyasi okur (ebooklib fallback)."""
    try:
        import ebooklib
        from ebooklib import epub
        from bs4 import BeautifulSoup
    except ImportError:
        raise ImportError(
            "EPUB destegi icin 'ebooklib' ve 'beautifulsoup4' gerekli: "
            "uv pip install ebooklib beautifulsoup4"
        )

    book = epub.read_epub(filepath)
    chapters = []
    idx = 0

    for item in book.get_items_of_type(ebooklib.ITEM_DOCUMENT):
        content = item.get_body_content()
        if content:
            soup = BeautifulSoup(content, "html.parser")
            text = soup.get_text(separator="\n", strip=True)
            if text.strip():
                idx += 1
                chapters.append(f"[Bolum {idx}]\n{text.strip()}")

    if not chapters:
        raise ValueError("EPUB dosyasindan metin cikarilmadi")
    return "\n\n---\n\n".join(chapters)


def _parse_json_python(filepath, ext):
    """JSON/JSONL dosyasi okur."""
    import json

    with open(filepath, "r", encoding="utf-8") as f:
        content = f.read()

    if ext == ".jsonl":
        entries = []
        for idx, line in enumerate(content.strip().splitlines(), 1):
            line = line.strip()
            if not line:
                continue
            try:
                obj = json.loads(line)
                text = _json_to_text(obj)
                if text:
                    entries.append(f"Kayit {idx}:\n{text}")
            except json.JSONDecodeError:
                continue
        if not entries:
            raise ValueError("JSONL dosyasindan veri cikarilmadi")
        return "\n\n".join(entries)
    else:
        data = json.loads(content)
        text = _json_to_text(data)
        if not text:
            raise ValueError("JSON dosyasindan veri cikarilmadi")
        return text


def _json_to_text(obj, indent=0):
    """JSON objesini okunabilir metne donusturur."""
    prefix = "  " * indent
    if isinstance(obj, dict):
        parts = []
        for key, value in obj.items():
            formatted_key = key.replace("_", " ").replace("-", " ").title()
            if isinstance(value, (dict, list)):
                sub = _json_to_text(value, indent + 1)
                if sub:
                    parts.append(f"{prefix}{formatted_key}:\n{sub}")
            else:
                parts.append(f"{prefix}{formatted_key}: {value}")
        return "\n".join(parts)
    elif isinstance(obj, list):
        if all(isinstance(item, (str, int, float, bool)) for item in obj):
            return f"{prefix}{', '.join(str(item) for item in obj)}"
        parts = []
        for idx, item in enumerate(obj, 1):
            sub = _json_to_text(item, indent + 1)
            if sub:
                parts.append(f"{prefix}[{idx}]\n{sub}")
        return "\n".join(parts)
    else:
        return f"{prefix}{obj}" if obj is not None else ""


# ── Belge Parse Yonlendirici ──────────────────────────────────

def _parse_file(filepath, ext):
    """Dosyayi uygun parser ile okur (Rust -> Python Fallback)."""
    if ext == ".pdf":
        return _parse_pdf(filepath)

    if RUST_PARSER_AVAILABLE:
        try:
            result = rust_parser.parse_document(filepath)
            return result.content
        except Exception as rust_err:
            print(f"  [UYARI] Rust parser basarisiz ({rust_err}), Python fallback deneniyor...")
            return _parse_with_python(filepath, ext)

    return _parse_with_python(filepath, ext)


# ── Belge Yukleme ─────────────────────────────────────────────

def load_documents(directory=None):
    """Belirtilen klasordeki tum desteklenen dosyalari okur."""
    if directory is None:
        directory = config.DOCUMENTS_DIR

    if not os.path.isdir(directory):
        raise FileNotFoundError(f"Belge klasoru bulunamadi: {directory}")

    supported_extensions = config.SUPPORTED_EXTENSIONS
    documents = []

    for ext in supported_extensions:
        pattern = os.path.join(directory, f"*{ext}")
        for filepath in sorted(glob.glob(pattern)):
            filename = os.path.basename(filepath)
            if filename.startswith("."):
                continue

            try:
                content = _parse_file(filepath, ext)
                if not content or not content.strip():
                    continue

                documents.append({
                    "filename": filename,
                    "filepath": filepath,
                    "content": content,
                })
                parser_type = "Rust" if (RUST_PARSER_AVAILABLE and ext != ".pdf") else "Python"
                print(f"  [OK] {filename} okundu ({len(content)} karakter) [{parser_type}]")

            except Exception as e:
                print(f"  [UYARI] {filename} okunamadi ({e}), atlaniyor.")
                continue

    if not documents:
        raise ValueError(f"'{directory}' klasorunde desteklenen belge bulunamadi.")

    return documents


# ── Chunk Kalite Filtresi ─────────────────────────────────────

def _is_quality_chunk(text):
    """Metin parcasinin yeterli kalitede olup olmadigini kontrol eder."""
    if not text or len(text.strip()) < 20:
        return False

    alnum_count = sum(1 for c in text if c.isalnum() or c.isspace())
    total_count = len(text)
    if total_count == 0:
        return False

    if (alnum_count / total_count) < 0.5:
        return False

    markup_indicators = ["[open]", "[close]", "[sep]", "<sep>", "end_of_msg", "role=\""]
    if sum(1 for ind in markup_indicators if ind in text.lower()) >= 2:
        return False

    if len(text.split()) < 3:
        return False

    return True


# ── Metin Parcalama (Chunking) ────────────────────────────────

def chunk_text(text, source_file, chunk_size=None, chunk_overlap=None):
    """Metni cumle sinirlarina duyarli olarak parcalara boler."""
    if not text or not text.strip():
        raise ValueError("Parcalanacak metin bos olamaz.")

    if chunk_size is None:
        chunk_size = config.CHUNK_SIZE
    if chunk_overlap is None:
        chunk_overlap = config.CHUNK_OVERLAP

    chunks = []
    start = 0
    chunk_index = 0

    while start < len(text):
        end = start + chunk_size

        if end < len(text):
            search_region = text[max(start, end - 100):end]
            last_period = -1
            for sep in [". ", ".\n", "? ", "?\n", "! ", "!\n", "\n\n"]:
                pos = search_region.rfind(sep)
                if pos > last_period:
                    last_period = pos

            if last_period > 0:
                end = max(start, end - 100) + last_period + 1

        chunk_text_piece = text[start:end].strip()

        if chunk_text_piece and _is_quality_chunk(chunk_text_piece):
            chunks.append({
                "text": chunk_text_piece,
                "source_file": source_file,
                "chunk_index": chunk_index,
            })
            chunk_index += 1

        start = end - chunk_overlap if end < len(text) else end

    return chunks


# ── Ana Ingestion Pipeline (Akilli Onbellek + Toplu Kayit) ────

def ingest_all(directory=None, clear_existing=False):
    """
    Akilli Ingestion Pipeline'ini calistirir:
    1. documents/ klasorundeki dosyalari hash kontrolu ile tarar.
    2. Degismeyen dosyalari atlar (0 sn).
    3. Sadece yeni veya degisen dosyalari parse edip embed eder.
    4. Qdrant'a toplu (bulk upsert) ile tek seferde kaydeder.

    Args:
        directory (str, optional): Belge klasoru.
        clear_existing (bool): True ise tum veritabani ve onbellegi sifirlayip bastan yukler.
    """
    start_total_time = time.time()

    print("\n" + "=" * 60)
    print("INGESTION PIPELINE (AKILLI ONBELLEK & BULK)")
    print(f"  Motor : {'Rust (yuksek performans)' if RUST_PARSER_AVAILABLE else 'Python (fallback)'}")
    print(f"  Mod   : {'Sıfırdan Tam Yükleme (--force)' if clear_existing else 'Akıllı Artımlı (Incremental Cache)'}")
    print("=" * 60)

    if directory is None:
        directory = config.DOCUMENTS_DIR

    if not os.path.isdir(directory):
        raise FileNotFoundError(f"Belge klasoru bulunamadi: {directory}")

    # 1. Veritabanini hazirla
    print("\n[1/4] Veritabani kontrol ediliyor...")
    if clear_existing:
        database.clear_db()
        if os.path.exists(CACHE_FILE):
            try:
                os.remove(CACHE_FILE)
            except Exception:
                pass
        print("  [OK] Mevcut veriler ve onbellek temizlendi.")
    database.init_db()

    cache = {} if clear_existing else _load_cache()
    db_sources = set(database.get_sources()) if not clear_existing else set()

    # 2. Dosya taramasi ve Hash analizi
    print(f"\n[2/4] Belgeler analiz ediliyor ({directory})...")
    current_files = {}
    for ext in config.SUPPORTED_EXTENSIONS:
        for filepath in glob.glob(os.path.join(directory, f"*{ext}")):
            filename = os.path.basename(filepath)
            if not filename.startswith("."):
                current_files[filename] = filepath

    files_to_process = []
    unchanged_files = []

    for filename, filepath in sorted(current_files.items()):
        fhash = _compute_file_hash(filepath)
        # Onbellekte varsa, hash ayniysa ve DB'de gercekten kaydi varsa atla
        if (
            not clear_existing
            and filename in cache
            and cache[filename].get("hash") == fhash
            and filename in db_sources
        ):
            unchanged_files.append(filename)
            print(f"  [ATLANDI] {filename} (Degisiklik yok, onbellekten korundu)")
        else:
            files_to_process.append((filename, filepath, fhash))
            # Eger eski kayit varsa guncelleme oncesi eski parcalari temizle
            if filename in db_sources:
                database.delete_file_chunks(filename)

    # Silinmis dosyalari DB ve onbellekten kaldir
    for cached_file in list(cache.keys()):
        if cached_file not in current_files:
            print(f"  [SILINDI] {cached_file} klasorden kaldirilmis, DB'den temizleniyor...")
            database.delete_file_chunks(cached_file)
            del cache[cached_file]

    if not files_to_process and unchanged_files:
        _save_cache(cache)
        total_chunks = database.get_chunk_count()
        print("\n" + "=" * 60)
        print("INGESTION TAMAMLANDI (TUM BELGELER GUNCELLI)")
        print(f"  Korunan Belge  : {len(unchanged_files)}")
        print(f"  Toplam DB Parca: {total_chunks}")
        print(f"  Gecen Sure     : {time.time() - start_total_time:.2f}s")
        print("=" * 60 + "\n")
        return {
            "documents_loaded": len(unchanged_files),
            "total_chunks": total_chunks,
            "files": unchanged_files,
        }

    # 3. Yeni/Degisen Belgeleri Oku ve Parcala
    print(f"\n[3/4] Yeni/Guncellenen belgeler okunuyor ve parcalaniyor ({len(files_to_process)} belge)...")
    new_chunks = []
    for filename, filepath, fhash in files_to_process:
        ext = os.path.splitext(filename)[1].lower()
        try:
            content = _parse_file(filepath, ext)
            if not content or not content.strip():
                continue

            doc_chunks = chunk_text(content, filename)
            new_chunks.extend(doc_chunks)
            cache[filename] = {
                "hash": fhash,
                "chunks_count": len(doc_chunks),
                "chars_count": len(content),
                "timestamp": time.time(),
            }
            parser_type = "Rust" if (RUST_PARSER_AVAILABLE and ext != ".pdf") else "Python"
            print(f"  [OK] {filename} -> {len(doc_chunks)} parca ({len(content)} kar.) [{parser_type}]")
        except Exception as e:
            print(f"  [UYARI] {filename} okunamadi ({e}), atlaniyor.")

    if not new_chunks and not unchanged_files:
        raise ValueError("Islenecek hicbir gecerli belge bulunamadi.")

    # 4. Toplu (Bulk) Embedding Uretimi ve Qdrant'a Kayit
    if new_chunks:
        total_new = len(new_chunks)
        print(f"\n[4/4] {total_new} yeni parca icin embedding uretiliyor ve kaydediliyor...")

        # 32'serli paketler halinde hem embed et hem aninda bulk upsert yap
        STEP_SIZE = 32
        for start_idx in range(0, total_new, STEP_SIZE):
            step_chunks = new_chunks[start_idx:start_idx + STEP_SIZE]
            step_texts = [c["text"] for c in step_chunks]

            # Embedding uret (optimize batch)
            step_embeddings = embeddings.get_embeddings_batch(step_texts)

            # Toplu kayit formatina cevir
            batch_to_insert = []
            for chunk, emb in zip(step_chunks, step_embeddings):
                batch_to_insert.append({
                    "text": chunk["text"],
                    "embedding": emb,
                    "source_file": chunk["source_file"],
                    "chunk_index": chunk["chunk_index"],
                })

            # Qdrant'a tek bulk istekle kaydet
            database.insert_chunks_batch(batch_to_insert)

            current_done = min(start_idx + STEP_SIZE, total_new)
            print(f"  [{current_done}/{total_new}] parca embed edildi ve kaydedildi...")

    # Onbellegi kaydet
    _save_cache(cache)

    total_chunks_in_db = database.get_chunk_count()
    all_active_files = list(current_files.keys())

    stats = {
        "documents_loaded": len(all_active_files),
        "new_documents_processed": len(files_to_process),
        "total_chunks": total_chunks_in_db,
        "files": all_active_files,
    }

    elapsed = time.time() - start_total_time
    record_ingestion_metrics(duration_sec=elapsed, documents_count=stats['documents_loaded'])

    print("\n" + "=" * 60)
    print("INGESTION TAMAMLANDI")
    print(f"  Toplam Belge     : {stats['documents_loaded']} (Yeni: {len(files_to_process)}, Onbellek: {len(unchanged_files)})")
    print(f"  Toplam DB Parca  : {stats['total_chunks']}")
    print(f"  Gecen Sure       : {elapsed:.2f}s")
    print("=" * 60 + "\n")

    return stats
