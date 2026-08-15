"""
Rust Parser + Yeni Dosya Formatlari Test Modulu.

Bu test modulu:
1. Rust parser'in Python'a dogru expose edildigini
2. Her format icin parse islemi calistigini
3. Python fallback mekanizmasinin duzgun calistigini
test eder.

Rust parser build edilmemisse, sadece Python fallback testleri calisir.
"""

import os
import sys
import json
import tempfile
import shutil

sys.path.insert(0, os.path.abspath(os.path.join(os.path.dirname(__file__), '..')))

from src import config

# ── Test Altyapisi ─────────────────────────────────────────────

PASSED = 0
FAILED = 0
SKIPPED = 0

def test(name, condition, detail=""):
    global PASSED, FAILED
    if condition:
        PASSED += 1
        print(f"  [OK] {name}")
    else:
        FAILED += 1
        print(f"  [FAIL] {name} -- {detail}")

def skip(name, reason):
    global SKIPPED
    SKIPPED += 1
    print(f"  [SKIP] {name} -- {reason}")


# ── Rust Parser Varlik Kontrolu ───────────────────────────────

print("=" * 60)
print("TEST: Rust Parser Varlik Kontrolu")
print("=" * 60)

try:
    import rust_parser
    RUST_AVAILABLE = True
    print("  [INFO] Rust parser MEVCUT")

    # Desteklenen uzantilar
    extensions = rust_parser.supported_extensions()
    test("supported_extensions() calisiyor", len(extensions) > 0)
    test("En az 10 format destekleniyor", len(extensions) >= 10,
         f"Bulunan: {len(extensions)}")
    print(f"  [INFO] Desteklenen: {', '.join(extensions)}")

except ImportError:
    RUST_AVAILABLE = False
    print("  [INFO] Rust parser MEVCUT DEGIL (sadece Python fallback test edilecek)")

print()


# ── Test Dosyalari Olustur ────────────────────────────────────

TEST_DIR = tempfile.mkdtemp(prefix="rag_test_")
print(f"  [INFO] Test dizini: {TEST_DIR}")

def create_test_file(filename, content, binary=False):
    """Test dosyasi olusturur."""
    path = os.path.join(TEST_DIR, filename)
    if binary:
        with open(path, "wb") as f:
            f.write(content)
    else:
        with open(path, "w", encoding="utf-8") as f:
            f.write(content)
    return path


# TXT dosyasi
txt_path = create_test_file("test.txt", "Bu bir test dosyasidir.\nIkinci satir.")

# MD dosyasi
md_path = create_test_file("test.md", "# Baslik\n\nBu bir markdown dosyasidir.\n\n## Alt Baslik\n\nIcerik burada.")

# CSV dosyasi
csv_path = create_test_file("test.csv", "Ad,Soyad,Yas\nAhmet,Yilmaz,30\nAyse,Kaya,25\nMehmet,Demir,35")

# TSV dosyasi
tsv_path = create_test_file("test.tsv", "Isim\tSehir\tYas\nAli\tIstanbul\t28\nVeli\tAnkara\t32")

# JSON dosyasi
json_data = {
    "name": "Test Projesi",
    "version": "1.0",
    "features": ["RAG", "Embedding", "LLM"],
    "author": {"name": "Ahmet", "email": "ahmet@test.com"}
}
json_path = create_test_file("test.json", json.dumps(json_data, ensure_ascii=False, indent=2))

# JSONL dosyasi
jsonl_lines = [
    json.dumps({"id": 1, "text": "Birinci kayit"}, ensure_ascii=False),
    json.dumps({"id": 2, "text": "Ikinci kayit"}, ensure_ascii=False),
    json.dumps({"id": 3, "text": "Ucuncu kayit"}, ensure_ascii=False),
]
jsonl_path = create_test_file("test.jsonl", "\n".join(jsonl_lines))

# HTML dosyasi
html_content = """<!DOCTYPE html>
<html>
<head><title>Test Sayfa</title></head>
<body>
<h1>Ana Baslik</h1>
<p>Bu bir test paragrafdir.</p>
<div>
    <h2>Alt Bolum</h2>
    <p>Ikinci paragraf.</p>
    <ul>
        <li>Madde 1</li>
        <li>Madde 2</li>
    </ul>
</div>
<script>var x = 1;</script>
<style>body { color: red; }</style>
</body>
</html>"""
html_path = create_test_file("test.html", html_content)

print()


# ── TEST 1: Duz Metin Dosyalari ──────────────────────────────

print("=" * 60)
print("TEST 1: Duz Metin (.txt, .md)")
print("=" * 60)

if RUST_AVAILABLE:
    result = rust_parser.parse_document(txt_path)
    test("TXT parse (Rust) -- content var", len(result.content) > 0)
    test("TXT parse (Rust) -- filename dogru", result.filename == "test.txt")
    test("TXT parse (Rust) -- format dogru", result.format == "txt")
    test("TXT parse (Rust) -- icerik dogru", "test dosyasidir" in result.content)

    result_md = rust_parser.parse_document(md_path)
    test("MD parse (Rust) -- content var", len(result_md.content) > 0)
    test("MD parse (Rust) -- icerik dogru", "markdown dosyasidir" in result_md.content)

# Python fallback
from src.ingestion import _parse_plaintext
py_result = _parse_plaintext(txt_path)
test("TXT parse (Python) -- content var", len(py_result) > 0)
test("TXT parse (Python) -- icerik dogru", "test dosyasidir" in py_result)

print()


# ── TEST 2: CSV/TSV ───────────────────────────────────────────

print("=" * 60)
print("TEST 2: Tablolu Veri (.csv, .tsv)")
print("=" * 60)

if RUST_AVAILABLE:
    result_csv = rust_parser.parse_document(csv_path)
    test("CSV parse (Rust) -- content var", len(result_csv.content) > 0)
    test("CSV parse (Rust) -- Ahmet iceriyor", "Ahmet" in result_csv.content)
    test("CSV parse (Rust) -- sutun basliklari var", "Ad" in result_csv.content)

    result_tsv = rust_parser.parse_document(tsv_path)
    test("TSV parse (Rust) -- content var", len(result_tsv.content) > 0)
    test("TSV parse (Rust) -- Istanbul iceriyor", "Istanbul" in result_tsv.content)

from src.ingestion import _parse_csv_python
py_csv = _parse_csv_python(csv_path, ".csv")
test("CSV parse (Python) -- content var", len(py_csv) > 0)
test("CSV parse (Python) -- Ahmet iceriyor", "Ahmet" in py_csv)

py_tsv = _parse_csv_python(tsv_path, ".tsv")
test("TSV parse (Python) -- content var", len(py_tsv) > 0)
test("TSV parse (Python) -- Istanbul iceriyor", "Istanbul" in py_tsv)

print()


# ── TEST 3: JSON/JSONL ───────────────────────────────────────

print("=" * 60)
print("TEST 3: Yapisal Veri (.json, .jsonl)")
print("=" * 60)

if RUST_AVAILABLE:
    result_json = rust_parser.parse_document(json_path)
    test("JSON parse (Rust) -- content var", len(result_json.content) > 0)
    test("JSON parse (Rust) -- proje adi var", "Test Projesi" in result_json.content)

    result_jsonl = rust_parser.parse_document(jsonl_path)
    test("JSONL parse (Rust) -- content var", len(result_jsonl.content) > 0)
    test("JSONL parse (Rust) -- kayitlar var", "kayit" in result_jsonl.content)

from src.ingestion import _parse_json_python
py_json = _parse_json_python(json_path, ".json")
test("JSON parse (Python) -- content var", len(py_json) > 0)
test("JSON parse (Python) -- proje adi var", "Test Projesi" in py_json)

py_jsonl = _parse_json_python(jsonl_path, ".jsonl")
test("JSONL parse (Python) -- content var", len(py_jsonl) > 0)
test("JSONL parse (Python) -- kayitlar var", "kayit" in py_jsonl)

print()


# ── TEST 4: HTML ──────────────────────────────────────────────

print("=" * 60)
print("TEST 4: HTML (.html)")
print("=" * 60)

if RUST_AVAILABLE:
    result_html = rust_parser.parse_document(html_path)
    test("HTML parse (Rust) -- content var", len(result_html.content) > 0)
    test("HTML parse (Rust) -- baslik var", "Ana Baslik" in result_html.content)
    test("HTML parse (Rust) -- script temizlendi", "var x" not in result_html.content)

from src.ingestion import _parse_html_python
py_html = _parse_html_python(html_path)
test("HTML parse (Python) -- content var", len(py_html) > 0)
test("HTML parse (Python) -- baslik var", "Ana Baslik" in py_html)
test("HTML parse (Python) -- script temizlendi", "var x" not in py_html)

print()


# ── TEST 5: DOCX (Python fallback) ───────────────────────────

print("=" * 60)
print("TEST 5: DOCX (.docx)")
print("=" * 60)

# DOCX icin gercek bir .docx dosyasi gerekiyor (binary ZIP format)
# python-docx ile test dosyasi olustur
try:
    from docx import Document
    doc = Document()
    doc.add_heading("Test Basligi", 0)
    doc.add_paragraph("Bu bir test paragrafdir.")
    doc.add_paragraph("Ikinci paragraf burada.")
    docx_path = os.path.join(TEST_DIR, "test.docx")
    doc.save(docx_path)

    if RUST_AVAILABLE:
        result_docx = rust_parser.parse_document(docx_path)
        test("DOCX parse (Rust) -- content var", len(result_docx.content) > 0)
        test("DOCX parse (Rust) -- baslik var", "Test Basligi" in result_docx.content)
        test("DOCX parse (Rust) -- paragraf var", "test paragrafdir" in result_docx.content)

    from src.ingestion import _parse_docx_python
    py_docx = _parse_docx_python(docx_path)
    test("DOCX parse (Python) -- content var", len(py_docx) > 0)
    test("DOCX parse (Python) -- baslik var", "Test Basligi" in py_docx)

except ImportError:
    skip("DOCX testleri", "python-docx yuklenmemis")

print()


# ── TEST 6: PPTX (Python fallback) ───────────────────────────

print("=" * 60)
print("TEST 6: PPTX (.pptx)")
print("=" * 60)

try:
    from pptx import Presentation
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Slide Basligi"
    slide.placeholders[1].text = "Slide icerik metni burada."
    pptx_path = os.path.join(TEST_DIR, "test.pptx")
    prs.save(pptx_path)

    if RUST_AVAILABLE:
        result_pptx = rust_parser.parse_document(pptx_path)
        test("PPTX parse (Rust) -- content var", len(result_pptx.content) > 0)
        test("PPTX parse (Rust) -- baslik var", "Slide Basligi" in result_pptx.content)

    from src.ingestion import _parse_pptx_python
    py_pptx = _parse_pptx_python(pptx_path)
    test("PPTX parse (Python) -- content var", len(py_pptx) > 0)
    test("PPTX parse (Python) -- baslik var", "Slide Basligi" in py_pptx)

except ImportError:
    skip("PPTX testleri", "python-pptx yuklenmemis")

print()


# ── TEST 7: XLSX (Python fallback) ───────────────────────────

print("=" * 60)
print("TEST 7: XLSX (.xlsx)")
print("=" * 60)

try:
    from openpyxl import Workbook
    wb = Workbook()
    ws = wb.active
    ws.title = "Veriler"
    ws.append(["Ad", "Soyad", "Sehir"])
    ws.append(["Ali", "Demir", "Istanbul"])
    ws.append(["Veli", "Kara", "Ankara"])
    xlsx_path = os.path.join(TEST_DIR, "test.xlsx")
    wb.save(xlsx_path)

    if RUST_AVAILABLE:
        result_xlsx = rust_parser.parse_document(xlsx_path)
        test("XLSX parse (Rust) -- content var", len(result_xlsx.content) > 0)
        test("XLSX parse (Rust) -- veri var", "Ali" in result_xlsx.content)

    from src.ingestion import _parse_xlsx_python
    py_xlsx = _parse_xlsx_python(xlsx_path)
    test("XLSX parse (Python) -- content var", len(py_xlsx) > 0)
    test("XLSX parse (Python) -- veri var", "Ali" in py_xlsx)

except ImportError:
    skip("XLSX testleri", "openpyxl yuklenmemis")

print()


# ── TEST 8: Batch Parse (Rust only) ──────────────────────────

if RUST_AVAILABLE:
    print("=" * 60)
    print("TEST 8: Batch Parse (Paralel)")
    print("=" * 60)

    batch_files = [txt_path, md_path, csv_path, json_path, html_path]
    batch_result = rust_parser.parse_documents_batch(batch_files)
    test("Batch parse -- toplam dogru", batch_result.total == len(batch_files))
    test("Batch parse -- hepsi basarili", len(batch_result.results) == len(batch_files))
    test("Batch parse -- hata yok", len(batch_result.errors) == 0)

    print()


# ── TEST 9: Hata Durumlari ────────────────────────────────────

print("=" * 60)
print("TEST 9: Hata Durumlari")
print("=" * 60)

if RUST_AVAILABLE:
    # Olmayan dosya
    try:
        rust_parser.parse_document("/olmayan/dosya.txt")
        test("Olmayan dosya hata veriyor (Rust)", False)
    except ValueError:
        test("Olmayan dosya hata veriyor (Rust)", True)

    # Desteklenmeyen format
    unsupported_path = create_test_file("test.xyz", "test")
    try:
        rust_parser.parse_document(unsupported_path)
        test("Desteklenmeyen format hata veriyor (Rust)", False)
    except ValueError:
        test("Desteklenmeyen format hata veriyor (Rust)", True)

# Python fallback hata
from src.ingestion import _parse_with_python
try:
    _parse_with_python("/olmayan/dosya.txt", ".txt")
    test("Olmayan dosya hata veriyor (Python)", False)
except Exception:
    test("Olmayan dosya hata veriyor (Python)", True)

try:
    _parse_with_python("test.xyz", ".xyz")
    test("Desteklenmeyen format hata veriyor (Python)", False)
except ValueError:
    test("Desteklenmeyen format hata veriyor (Python)", True)

print()


# ── TEST 10: load_documents() Entegrasyonu ────────────────────

print("=" * 60)
print("TEST 10: load_documents() ile entegrasyon")
print("=" * 60)

from src.ingestion import load_documents

# Test dizinine birden fazla format koy
multi_dir = os.path.join(TEST_DIR, "multi_format")
os.makedirs(multi_dir, exist_ok=True)
shutil.copy(txt_path, multi_dir)
shutil.copy(csv_path, multi_dir)
shutil.copy(json_path, multi_dir)
shutil.copy(html_path, multi_dir)

docs = load_documents(multi_dir)
test("load_documents -- birden fazla format yuklendi", len(docs) >= 4,
     f"Yuklenen: {len(docs)}")

formats_found = set(os.path.splitext(d["filename"])[1] for d in docs)
test("load_documents -- .txt yuklendi", ".txt" in formats_found)
test("load_documents -- .csv yuklendi", ".csv" in formats_found)
test("load_documents -- .json yuklendi", ".json" in formats_found)
test("load_documents -- .html yuklendi", ".html" in formats_found)

for doc in docs:
    test(f"load_documents -- {doc['filename']} bos degil",
         len(doc["content"]) > 0)

print()


# ── TEST 11: Config.SUPPORTED_EXTENSIONS ──────────────────────

print("=" * 60)
print("TEST 11: Config")
print("=" * 60)

test("SUPPORTED_EXTENSIONS tanimli", hasattr(config, "SUPPORTED_EXTENSIONS"))
test("En az 12 format var", len(config.SUPPORTED_EXTENSIONS) >= 12,
     f"Bulunan: {len(config.SUPPORTED_EXTENSIONS)}")
test(".docx destekleniyor", ".docx" in config.SUPPORTED_EXTENSIONS)
test(".pptx destekleniyor", ".pptx" in config.SUPPORTED_EXTENSIONS)
test(".html destekleniyor", ".html" in config.SUPPORTED_EXTENSIONS)
test(".csv destekleniyor", ".csv" in config.SUPPORTED_EXTENSIONS)
test(".xlsx destekleniyor", ".xlsx" in config.SUPPORTED_EXTENSIONS)
test(".epub destekleniyor", ".epub" in config.SUPPORTED_EXTENSIONS)
test(".json destekleniyor", ".json" in config.SUPPORTED_EXTENSIONS)
test(".jsonl destekleniyor", ".jsonl" in config.SUPPORTED_EXTENSIONS)

print()


# ── Temizlik ──────────────────────────────────────────────────
shutil.rmtree(TEST_DIR, ignore_errors=True)


# ── SONUC ─────────────────────────────────────────────────────
print("=" * 60)
total = PASSED + FAILED + SKIPPED
print(f"SONUC: {PASSED}/{total} GECTI, {FAILED} BASARISIZ, {SKIPPED} ATLANDI")
if RUST_AVAILABLE:
    print("  Motor: Rust parser + Python fallback")
else:
    print("  Motor: Python fallback (Rust parser build edilmemis)")
if FAILED > 0:
    print("  BAZI TESTLER BASARISIZ!")
else:
    print("  TUMU BASARILI!")
print("=" * 60)

sys.exit(0 if FAILED == 0 else 1)
